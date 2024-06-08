#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import pandas as pd
from googletrans import Translator
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfFileReader, PdfFileWriter
import io

def load_dictionaries(csv_files):
    translation_dict = {}
    for file in csv_files:
        df = pd.read_csv(file)
        for index, row in df.iterrows():
            indonesian_word = row['indonesian']
            english_word = row['english']
            translation_dict[indonesian_word] = english_word
            translation_dict[english_word] = indonesian_word
    return translation_dict

def translate_text(text, translation_dict):
    words = text.split()
    translated_words = [translation_dict.get(word, word) for word in words]
    return ' '.join(translated_words)

def translate_with_google(text, target_lang, translation_dict):
    text = translate_text(text, translation_dict)
    translator = Translator()
    detected_lang = translator.detect(text).lang
    translated = translator.translate(text, src=detected_lang, dest=target_lang)
    return translated.text

def translate_word(doc, target_lang, translation_dict):
    for paragraph in doc.paragraphs:
        paragraph.text = translate_with_google(paragraph.text, target_lang, translation_dict)
    return doc

def translate_excel(wb, target_lang, translation_dict):
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value:
                    cell.value = translate_with_google(cell.value, target_lang, translation_dict)
    return wb

def translate_pptx(prs, target_lang, translation_dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = translate_with_google(run.text, target_lang, translation_dict)
    return prs

def translate_pdf(pdf_reader, target_lang, translation_dict):
    pdf_writer = PdfFileWriter()
    for page_num in range(pdf_reader.numPages):
        page = pdf_reader.getPage(page_num)
        page_content = page.extract_text()
        translated_content = translate_with_google(page_content, target_lang, translation_dict)
        page.merge_page(PdfFileReader(io.BytesIO(translated_content.encode('utf-8'))).getPage(0))
        pdf_writer.addPage(page)
    return pdf_writer

def save_translated_file(translated_file, file_type):
    if file_type == "docx":
        translated_file.save("translated_document.docx")
    elif file_type == "xlsx":
        translated_file.save("translated_spreadsheet.xlsx")
    elif file_type == "pptx":
        translated_file.save("translated_presentation.pptx")
    elif file_type == "pdf":
        with open("translated_document.pdf", "wb") as f:
            translated_file.write(f)

